import pptx 
import os
from format import create_json_file

presentation = pptx.Presentation(os.path.join("slide", "template.pptx"))

data = create_json_file()

def full_slides(data):

    title = data["title_page"][0]["title"]
    edit_title_slide(title)

    num_of_slides = len(data["Slides"])

    slide_number = 0
    while(slide_number < num_of_slides):

        slide_data = data["Slides"][slide_number]
        edit_slides(slide_number, slide_data)
        slide_number+=1


def edit_slides(slide_number, slide_data):
    slide = presentation.slides[slide_number + 1]
    count = 0
    for shape in slide.shapes:
        if shape.has_text_frame and shape.text:
            count += 1
            # If the textbox is the heading 
            if count == 1:
                text_frame = shape.text_frame

                first_paragraph = text_frame.paragraphs[0]
                first_run = first_paragraph.runs[0] if first_paragraph.runs else first_paragraph.add_run()

                # Preserve formatting of first run
                font = first_run.font
                font_name = font.name
                font_size = font.size
                font_bold = font.bold
                font_italic = font.italic
                font_underline = font.underline
                font_colour = font.color.rgb

                # Clear existing text in frame and apply new text with preserved formatting
                text_frame.clear()
                new_run = text_frame.paragraphs[0].add_run()
                new_run.text = slide_data["slide_title"]

                new_run.font.name = font_name
                new_run.font.size = font_size
                new_run.font.bold = font_bold
                new_run.font.italic = font_italic
                new_run.font.underline = font_underline
                new_run.font.color.rgb = font_colour
            
            # If the textbox is the bulletpoint section
            if count == 2:  
                text_frame = shape.text_frame

                first_paragraph = text_frame.paragraphs[0]
                first_run = first_paragraph.runs[0] if first_paragraph.runs else first_paragraph.add_run()

                # Preserve formatting of first run
                font = first_run.font
                font_name = font.name
                font_size = font.size
                font_bold = font.bold
                font_italic = font.italic
                font_underline = font.underline
                font_colour = font.color.rgb

                # Clear existing text in frame and apply new text with preserved formatting
                text_frame.clear()
                new_run = text_frame.paragraphs[0].add_run()

                bulletpoints = slide_data["bullet_points"][0]

                for idx, key in enumerate(bulletpoints, start=1):
                    paragraph = text_frame.add_paragraph()
                    run = paragraph.add_run()
                    run.text = f"{idx}. {bulletpoints[key]}"

                    # Apply preserved formatting
                    run.font.name = font_name
                    run.font.size = font_size
                    run.font.bold = font_bold
                    run.font.italic = font_italic
                    run.font.underline = font_underline
                    run.font.color.rgb = font_colour
                return


def edit_title_slide(title):
    slide = presentation.slides[0]
    count = 0
    for shape in slide.shapes:
        if shape.has_text_frame and shape.text:
            count += 1
            # If the textbox is the subheading text box
            if count == 1:
                text_frame = shape.text_frame

                first_paragraph = text_frame.paragraphs[0]
                first_run = first_paragraph.runs[0] if first_paragraph.runs else first_paragraph.add_run()

                # Preserve formatting of first run
                font = first_run.font
                font_name = font.name
                font_size = font.size
                font_bold = font.bold
                font_italic = font.italic
                font_underline = font.underline
                font_colour = font.color.rgb

                # Clear existing text in frame and apply new text with preserved formatting
                text_frame.clear()
                new_run = text_frame.paragraphs[0].add_run()
                new_run.text = "beta 1.0"

                new_run.font.name = font_name
                new_run.font.size = font_size
                new_run.font.bold = font_bold
                new_run.font.italic = font_italic
                new_run.font.underline = font_underline
                new_run.font.color.rgb = font_colour
            # If the textbox is the heading textbox     
            if count == 2:
                text_frame = shape.text_frame

                first_paragraph = text_frame.paragraphs[0]
                first_run = first_paragraph.runs[0] if first_paragraph.runs else first_paragraph.add_run()

                # Preserve formatting of first run
                font = first_run.font
                font_name = font.name
                font_size = font.size
                font_bold = font.bold
                font_italic = font.italic
                font_underline = font.underline
                font_colour = font.color.rgb

                # Clear existing text in frame and apply new text with preserved formatting
                text_frame.clear()
                new_run = text_frame.paragraphs[0].add_run()
                new_run.text = title

                new_run.font.name = font_name
                new_run.font.size = font_size
                new_run.font.bold = font_bold
                new_run.font.italic = font_italic
                new_run.font.underline = font_underline
                new_run.font.color.rgb = font_colour
                return

def list_text_boxes(presentation, slide_num):
    slide = presentation.slides[slide_num - 1]
    text_boxes = []
    for shape in slide.shapes:
        if shape.has_text_frame and shape.text:
            text_boxes.append(shape.text)
    return text_boxes

def update_text_of_textbox(presentation, slide_num, textbox_id, new_text):
    slide = presentation.slides[(slide_num - 1)]
    count = 0
    for shape in slide.shapes:
        if shape.has_text_frame and shape.text:
            count += 1
            if count == textbox_id:
                text_frame = shape.text_frame
                first_paragraph = text_frame.paragraphs[0]
                first_run = first_paragraph.runs[0] if first_paragraph.runs else first_paragraph.add_run()

                # Preserve formatting of first run
                font = first_run.font
                font_name = font.name
                font_size = font.size
                font_bold = font.bold
                font_italic = font.italic
                font_underline = font.underline
                font_colour = font.color.rgb

                # Clear existing text in frame and apply new text with preserved formatting
                text_frame.clear()
                new_run = text_frame.paragraphs[0].add_run()
                new_run.text = new_text

                new_run.font.name = font_name
                new_run.font.size = font_size
                new_run.font.bold = font_bold
                new_run.font.italic = font_italic
                new_run.font.underline = font_underline
                new_run.font.color.rgb = font_colour
                return


for idx, text in enumerate(list_text_boxes(presentation=presentation, slide_num=1), 1):
    print(f"Text Box {idx}: {text}")

title = "The History of Porsche"
subheading = "A course that goes into the history of Porsche and how it was founded"

full_slides(data)

output_path = "new" + "_template.pptx"
presentation.save(output_path)